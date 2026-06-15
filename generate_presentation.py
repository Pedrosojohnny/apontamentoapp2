import os
import copy
import shutil
from pptx import Presentation

def set_text_preserve_format(shape, text):
    """
    Substitui o texto de uma forma (shape) mantendo a formatação
    de caractere do primeiro run (fonte, tamanho, cor, negrito, etc.).
    """
    if not shape.has_text_frame:
        return
    tf = shape.text_frame
    
    if tf.paragraphs:
        p = tf.paragraphs[0]
        
        # Remove todos os parágrafos extras (deixando apenas o primeiro)
        # Acessa o XML diretamente para remover os elementos filhos 'p'
        p_elements = list(tf._element.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}p'))
        for p_el in p_elements[1:]:
            tf._element.remove(p_el)
            
        # Agora limpa as runs adicionais no primeiro parágrafo
        if p.runs:
            p.runs[0].text = text
            r_elements = list(p._element.findall('{http://schemas.openxmlformats.org/drawingml/2006/main}r'))
            for r_el in r_elements[1:]:
                p._element.remove(r_el)
        else:
            p.text = text
    else:
        tf.text = text

def clone_selected_shapes(src_slide, dest_slide, names_to_keep):
    """
    Clona apenas as formas especificadas de um slide de origem para um slide de destino,
    mantendo a ordem original de camadas (z-order).
    """
    shapes_to_clone = []
    for idx, shape in enumerate(src_slide.shapes):
        if shape.name in names_to_keep:
            shapes_to_clone.append((idx, shape))
            
    # Ordena por índice de camada original
    shapes_to_clone.sort(key=lambda x: x[0])
    
    cloned_shapes = {}
    for _, shape in shapes_to_clone:
        new_el = copy.deepcopy(shape.element)
        dest_slide.shapes._spTree.append(new_el)
        cloned_shape = dest_slide.shapes[-1]
        cloned_shapes[shape.name] = cloned_shape
        
    return cloned_shapes

def generate_presentation(src_path, dest_path):
    """
    Gera a nova apresentação do zero, usando apenas o design visual,
    backgrounds, logos e formatos de caixas de texto do template.
    """
    if not os.path.exists(src_path):
        print(f"Erro: Arquivo de template '{src_path}' não encontrado.")
        return False
        
    # Abre a apresentação de exemplo para usar como fonte dos slides master e layouts
    prs = Presentation(src_path)
    
    # Guarda os slides de template originais em memória
    templates = list(prs.slides)
    print(f"Carregados {len(templates)} slides de template.")
    
    # Deleta todos os slides existentes da apresentação principal
    id_list = prs.slides._sldIdLst
    for idx in range(len(id_list) - 1, -1, -1):
        rId = id_list[idx].rId
        prs.part.drop_rel(rId)
        del id_list[idx]
    print("Todos os slides originais foram limpos da apresentação destino.")

    # Configuração de cada novo slide a ser gerado
    # Mapeia qual layout usar (0 = DEFAULT, 1 = 6_Layout Personalizado) e quais elementos clonar
    slides_config = [
        {
            # Slide 1: Capa
            "layout_idx": 0,
            "src_slide": templates[0],
            "keep": ["Image 0", "Imagem 7", "Text 0", "Text 1", "CaixaDeTexto 8"],
            "texts": {
                "Text 0": "SISTEMA DE APONTAMENTO",
                "Text 1": "Automação e Registro de Produção de Painéis",
                "CaixaDeTexto 8": "PCP / João Vitor Pedroso"
            }
        },
        {
            # Slide 2: Visão Geral
            "layout_idx": 0,
            "src_slide": templates[1],
            "keep": ["Image 0", "Shape 0", "Imagem 31", "Text 1", "Text 7", "Text 8", "Text 9", "Text 10", 
                     "Imagem 19", "Imagem 20", "Imagem 21", "Imagem 22"],
            "texts": {
                "Text 1": "Visão Geral do Sistema",
                "Text 7": "Apontamento Digital: Substitui o registro manual e o retrabalho de digitação no PCP.",
                "Text 8": "Operação por Tablets: Coleta de tempos usando câmera para leitura de código de barras.",
                "Text 9": "Saneamento Automático: Higienização de IDs, OMs e Seriais eliminando zeros à esquerda.",
                "Text 10": "Histórico Offline: Armazenamento local (IndexedDB) para garantir operação ininterrupta."
            }
        },
        {
            # Slide 3: Tecnologia e Arquitetura
            "layout_idx": 0,
            "src_slide": templates[2],
            "keep": ["Image 0", "Shape 0", "Imagem 26", "Text 1", "Text 4", "Text 5", "Text 8", "Text 9", 
                     "Text 12", "Text 13", "Text 16", "Text 17", "Imagem 29", "Imagem 30", "Imagem 31", "Imagem 32"],
            "texts": {
                "Text 1": "Tecnologia e Arquitetura",
                "Text 4": "Operação Offline",
                "Text 5": "Armazenamento local via IndexedDB evita perda de registros em oscilações de rede.",
                "Text 8": "Validação Rápida",
                "Text 9": "Garante sequenciamento do roteiro e trava de quantidade na própria máquina.",
                "Text 12": "Leitura Óptica",
                "Text 13": "Biblioteca html5-qrcode realiza a leitura rápida de Code 128 em baixa luz.",
                "Text 16": "Integração PCP",
                "Text 17": "Importação de base de roteiros e exportação consolidada de dados por Excel (SheetJS)."
            }
        },
        {
            # Slide 4: Fluxo de Trabalho (baseado na estrutura do slide 2)
            "layout_idx": 0,
            "src_slide": templates[1],
            "keep": ["Image 0", "Shape 0", "Imagem 31", "Text 1", "Text 7", "Text 8", "Text 9", "Text 10", 
                     "Imagem 19", "Imagem 20", "Imagem 21", "Imagem 22"],
            "texts": {
                "Text 1": "Fluxo de Trabalho do Operador",
                "Text 7": "1. Identificação: Operador faz login digitando o número de ID ou escaneando o crachá.",
                "Text 8": "2. Configuração: Seleção da Ordem de Montagem (OM) e a etapa que será realizada.",
                "Text 9": "3. Início da Unidade: Primeiro bip na etiqueta de número de série inicia a cronometragem.",
                "Text 10": "4. Finalização: Segundo bip no mesmo número de série encerra o tempo e salva o registro."
            }
        },
        {
            # Slide 5: Regras de Negócio e Segurança
            "layout_idx": 0,
            "src_slide": templates[4],
            "keep": ["Image 0", "Shape 0", "Imagem 31", "Text 1", "Text 9", "Text 10", "Text 13", "Text 14", 
                     "Text 17", "Text 18", "Text 21", "Text 22", "Imagem 26", "Imagem 27", "Imagem 28", "Imagem 29"],
            "texts": {
                "Text 1": "Regras de Negócio e Segurança",
                "Text 9": "Gestão de Pausas",
                "Text 10": "Controle automático de pausas do operador, descontando paradas do tempo líquido.",
                "Text 13": "Trava de Sequência",
                "Text 14": "Impede o início de uma etapa se a peça não passou pela operação anterior obrigatória.",
                "Text 17": "Trava de Quantidade",
                "Text 18": "Bloqueia novas leituras quando a meta planejada na OM for atingida.",
                "Text 21": "Retomada de Peça",
                "Text 22": "Opção 'Finalizar Operação' permite registrar peças que iniciaram em turnos anteriores."
            }
        },
        {
            # Slide 6: Painel Administrativo PCP
            "layout_idx": 0,
            "src_slide": templates[5],
            "keep": ["Image 0", "Shape 0", "Imagem 29", "Text 1", "Text 2", "Text 4", "Text 6", "Text 8", "Text 10",
                     "Imagem 25", "Imagem 26", "Imagem 27", "Imagem 28"],
            "texts": {
                "Text 1": "Painel de Administração PCP",
                "Text 2": "Configurações Gerais e Gestão de Apontamentos",
                "Text 4": "Acesso Restrito: Área protegida por senha padrão (PCP2000) para controle seguro.",
                "Text 6": "Sincronização: Upload simples de planilhas de Ordens e Roteiros com mapeamento inteligente.",
                "Text 8": "Edição de Dados: Permite correção manual de horários de início/fim e tempos de pausa.",
                "Text 10": "Exportação Excel: Relatório pronto com tempos consolidados e limpeza segura da base."
            }
        },
        {
            # Slide 7: Contatos Bauer (usando layout personalizado 6_Layout Personalizado)
            "layout_idx": 1,
            "src_slide": templates[6],
            "keep": ["T\u00edtulo 1", "Espa\u00e7o Reservado para Texto 3", "Espa\u00e7o Reservado para Texto 4", 
                     "Espa\u00e7o Reservado para Texto 5", "Espa\u00e7o Reservado para Texto 6", "Espa\u00e7o Reservado para Texto 7", 
                     "Espa\u00e7o Reservado para Texto 8", "Espa\u00e7o Reservado para Texto 9", "Espa\u00e7o Reservado para Texto 10"],
            "texts": {} # Mantém informações de contato intactas
        }
    ]

    for s_idx, cfg in enumerate(slides_config):
        print(f"Gerando Slide {s_idx + 1}...")
        
        # Adiciona slide com o respectivo layout
        layout = prs.slide_layouts[cfg["layout_idx"]]
        new_slide = prs.slides.add_slide(layout)
        
        # Clona os elementos permitidos
        cloned_shapes = clone_selected_shapes(cfg["src_slide"], new_slide, cfg["keep"])
        
        # Atualiza os textos mantendo formatação
        for name, text in cfg["texts"].items():
            if name in cloned_shapes:
                set_text_preserve_format(cloned_shapes[name], text)
                
    # Salva o arquivo final
    prs.save(dest_path)
    print(f"\nApresentação gerada com sucesso em '{dest_path}'!")
    return True

if __name__ == "__main__":
    src_file = "Exemplo apresentação.pptx"
    dest_file = "Apresentação Sistema de Apontamento.pptx"
    generate_presentation(src_file, dest_file)
